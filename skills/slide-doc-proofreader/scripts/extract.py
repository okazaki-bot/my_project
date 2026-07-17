#!/usr/bin/env python3
"""
スライド/ドキュメントからテキストを「位置情報つき」で抽出する。

校正（誤字脱字・不適切表現）のチェックは、抽出したテキストを人/モデルが
読んで判断する作業なので、このスクリプトの役割は「どこに何が書いてあるか」を
正確に・機械可読な形で取り出すことに徹する。

対応形式: .pptx / .docx / .pdf

出力(JSON, stdout):
{
  "type": "pptx|docx|pdf",
  "path": "...",
  "units": [
    {
      "index": 1,                 # スライド番号 / ページ番号（1始まり）
      "label": "スライド1",
      "full_text": "...",         # 重複検出・通読用にユニット全体を連結したテキスト
      "blocks": [
        {"id": "u1-b0", "kind": "title|body|table|text", "text": "..."}
      ]
    }
  ]
}

block.id はそのまま apply_fixes.py の置換対象指定にも使える。
"""
import json
import sys
import os


def extract_pptx(path):
    from pptx import Presentation
    from pptx.util import Emu  # noqa
    prs = Presentation(path)
    units = []
    for si, slide in enumerate(prs.slides, start=1):
        blocks = []
        bi = 0
        for shape in slide.shapes:
            kind = "body"
            try:
                if shape.is_placeholder and shape.placeholder_format is not None:
                    ph = shape.placeholder_format.type
                    if ph is not None and "TITLE" in str(ph):
                        kind = "title"
            except Exception:
                pass
            # 通常のテキスト枠
            if shape.has_text_frame:
                text = "\n".join(p.text for p in shape.text_frame.paragraphs)
                if text.strip():
                    blocks.append({"id": f"u{si}-b{bi}", "kind": kind, "text": text})
                    bi += 1
            # 表
            if shape.has_table:
                for r, row in enumerate(shape.table.rows):
                    for c, cell in enumerate(row.cells):
                        ct = cell.text
                        if ct.strip():
                            blocks.append({
                                "id": f"u{si}-b{bi}",
                                "kind": "table",
                                "cell": [r, c],
                                "text": ct,
                            })
                            bi += 1
        # ノート（スピーカーノート）も拾う
        try:
            if slide.has_notes_slide:
                nt = slide.notes_slide.notes_text_frame.text
                if nt.strip():
                    blocks.append({"id": f"u{si}-b{bi}", "kind": "notes", "text": nt})
                    bi += 1
        except Exception:
            pass
        full = "\n".join(b["text"] for b in blocks if b["kind"] != "notes")
        units.append({
            "index": si,
            "label": f"スライド{si}",
            "full_text": full,
            "blocks": blocks,
        })
    return {"type": "pptx", "path": path, "units": units}


def extract_docx(path):
    from docx import Document
    doc = Document(path)
    blocks = []
    bi = 0
    for p in doc.paragraphs:
        if p.text.strip():
            style = p.style.name if p.style is not None else ""
            kind = "title" if "Heading" in style or "Title" in style else "body"
            blocks.append({"id": f"p-b{bi}", "kind": kind, "style": style, "text": p.text})
            bi += 1
    for ti, table in enumerate(doc.tables):
        for r, row in enumerate(table.rows):
            for c, cell in enumerate(row.cells):
                ct = cell.text
                if ct.strip():
                    blocks.append({
                        "id": f"p-b{bi}",
                        "kind": "table",
                        "cell": [ti, r, c],
                        "text": ct,
                    })
                    bi += 1
    # docx はページ概念が曖昧なので、ドキュメント全体を 1 ユニットとして扱い、
    # 段落を blocks として並べる。重複検出は段落単位で別途行う。
    full = "\n".join(b["text"] for b in blocks)
    units = [{"index": 1, "label": "本文", "full_text": full, "blocks": blocks}]
    return {"type": "docx", "path": path, "units": units}


def extract_pdf(path):
    from pypdf import PdfReader
    reader = PdfReader(path)
    units = []
    for pi, page in enumerate(reader.pages, start=1):
        text = page.extract_text() or ""
        blocks = []
        if text.strip():
            blocks.append({"id": f"u{pi}-b0", "kind": "text", "text": text})
        units.append({
            "index": pi,
            "label": f"ページ{pi}",
            "full_text": text,
            "blocks": blocks,
        })
    return {"type": "pdf", "path": path, "units": units}


def main():
    if len(sys.argv) < 2:
        print("usage: extract.py <file.pptx|.docx|.pdf>", file=sys.stderr)
        sys.exit(2)
    path = sys.argv[1]
    ext = os.path.splitext(path)[1].lower()
    if ext == ".pptx":
        data = extract_pptx(path)
    elif ext == ".docx":
        data = extract_docx(path)
    elif ext == ".pdf":
        data = extract_pdf(path)
    else:
        print(f"未対応の形式です: {ext}（.pptx/.docx/.pdf のみ対応）", file=sys.stderr)
        sys.exit(2)
    print(json.dumps(data, ensure_ascii=False, indent=2))


if __name__ == "__main__":
    main()
