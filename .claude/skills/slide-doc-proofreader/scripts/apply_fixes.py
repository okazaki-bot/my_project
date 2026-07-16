#!/usr/bin/env python3
"""
指摘に対するユーザーの指示にもとづいて、スライド/ドキュメントを実際に修正する。

安全のため、修正前に必ず元ファイルのバックアップを作ってから、元ファイルを
上書き保存する（リンクや参照を壊さないため、出力先は元のパスのまま）。

使い方:
    apply_fixes.py <file.pptx|.docx|.pdf> <edits.json>

edits.json は編集指示の配列:
[
  {"action": "replace", "find": "誤った語", "replace": "正しい語", "unit": 3},
  {"action": "replace", "find": "全体置換語", "replace": "修正語"},
  {"action": "delete_unit", "unit": 7}
]

- action "replace": find→replace の文字列置換。
    "unit" を指定するとそのスライド/ページ内だけに限定（指定なしは全体）。
    docx は unit 概念が無いので常に本文全体が対象。
- action "delete_unit": スライド/ページを丸ごと削除（重複スライドの除去用）。

PDF はテキストの書き換えが安全に行えないため "replace" は非対応
（スキップして報告）。"delete_unit"（ページ削除）は対応する。

出力(JSON, stdout): 適用結果のレポート。
"""
import json
import sys
import os
import shutil
from datetime import datetime


def backup(path):
    base, ext = os.path.splitext(path)
    ts = datetime.now().strftime("%Y%m%d-%H%M%S")
    bak = f"{base}.bak-{ts}{ext}"
    shutil.copy2(path, bak)
    return bak


def _replace_in_paragraphs(paragraphs, find, replace):
    """run の書式を保ちつつ段落テキストを置換。
    段落内の最初の run にまとめ、残りを空にする方式。"""
    n = 0
    for para in paragraphs:
        runs = para.runs
        if not runs:
            continue
        full = "".join(r.text for r in runs)
        if find in full:
            n += full.count(find)
            runs[0].text = full.replace(find, replace)
            for r in runs[1:]:
                r.text = ""
    return n


# ---------- PPTX ----------
def fix_pptx(path, edits):
    from pptx import Presentation
    prs = Presentation(path)
    slides = list(prs.slides)
    report = []

    # 削除はインデックスがずれないよう後ろから処理
    deletes = sorted([e["unit"] for e in edits if e["action"] == "delete_unit"],
                     reverse=True)
    replaces = [e for e in edits if e["action"] == "replace"]

    for e in replaces:
        find, rep = e["find"], e["replace"]
        target_units = [e["unit"]] if e.get("unit") else range(1, len(slides) + 1)
        cnt = 0
        for ui in target_units:
            slide = slides[ui - 1]
            for shape in slide.shapes:
                if shape.has_text_frame:
                    cnt += _replace_in_paragraphs(shape.text_frame.paragraphs, find, rep)
                if shape.has_table:
                    for row in shape.table.rows:
                        for cell in row.cells:
                            cnt += _replace_in_paragraphs(cell.text_frame.paragraphs, find, rep)
        report.append({"action": "replace", "find": find, "replace": rep,
                       "unit": e.get("unit", "all"), "replaced": cnt})

    if deletes:
        sld_id_lst = prs.slides._sldIdLst
        sld_ids = list(sld_id_lst)
        for ui in deletes:
            sld_id_lst.remove(sld_ids[ui - 1])
            report.append({"action": "delete_unit", "unit": ui, "deleted": True})

    prs.save(path)
    return report


# ---------- DOCX ----------
def fix_docx(path, edits):
    from docx import Document
    doc = Document(path)
    report = []
    for e in edits:
        if e["action"] != "replace":
            report.append({"action": e["action"], "skipped": "docxではunit削除は未対応"})
            continue
        find, rep = e["find"], e["replace"]
        cnt = _replace_in_paragraphs(doc.paragraphs, find, rep)
        for table in doc.tables:
            for row in table.rows:
                for cell in row.cells:
                    cnt += _replace_in_paragraphs(cell.paragraphs, find, rep)
        report.append({"action": "replace", "find": find, "replace": rep, "replaced": cnt})
    doc.save(path)
    return report


# ---------- PDF ----------
def fix_pdf(path, edits):
    from pypdf import PdfReader, PdfWriter
    report = []
    deletes = sorted([e["unit"] for e in edits if e["action"] == "delete_unit"])
    for e in edits:
        if e["action"] == "replace":
            report.append({"action": "replace", "find": e["find"],
                           "skipped": "PDFのテキスト書き換えは非対応。元データ(pptx/docx等)で修正してください"})
    if deletes:
        reader = PdfReader(path)
        writer = PdfWriter()
        for i, page in enumerate(reader.pages, start=1):
            if i in deletes:
                report.append({"action": "delete_unit", "unit": i, "deleted": True})
                continue
            writer.add_page(page)
        with open(path, "wb") as f:
            writer.write(f)
    return report


def main():
    if len(sys.argv) < 3:
        print("usage: apply_fixes.py <file> <edits.json>", file=sys.stderr)
        sys.exit(2)
    path, edits_path = sys.argv[1], sys.argv[2]
    with open(edits_path, encoding="utf-8") as f:
        edits = json.load(f)
    if not isinstance(edits, list):
        edits = edits.get("edits", [])

    bak = backup(path)
    ext = os.path.splitext(path)[1].lower()
    if ext == ".pptx":
        report = fix_pptx(path, edits)
    elif ext == ".docx":
        report = fix_docx(path, edits)
    elif ext == ".pdf":
        report = fix_pdf(path, edits)
    else:
        print(f"未対応の形式です: {ext}", file=sys.stderr)
        sys.exit(2)

    print(json.dumps({
        "file": path,
        "backup": bak,
        "applied": report,
    }, ensure_ascii=False, indent=2))


if __name__ == "__main__":
    main()
